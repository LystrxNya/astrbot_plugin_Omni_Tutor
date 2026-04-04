import sys
import subprocess
import warnings
warnings.filterwarnings("ignore", category=UserWarning) # 忽略 pandas 的某些底层警告

# ================= 自动依赖安装逻辑 =================
try:
    import docx
    import pptx
    import pandas as pd
    import tabulate
except ImportError:
    print("[OmniTutor] 正在自动安装办公文档解析依赖 (约需几十秒)，请稍候...")
    subprocess.check_call([
        sys.executable, "-m", "pip", "install", 
        "python-docx", "python-pptx", "pandas", "openpyxl", "tabulate",
        "-i", "https://pypi.tuna.tsinghua.edu.cn/simple"
    ])
    import docx
    import pptx
    import pandas as pd

def parse_office_file(file_path: str, ext: str) -> str:
    """
    通用办公文档解析器，将各种格式转换为 Markdown 纯文本
    """
    text = ""
    try:
        if ext in ['txt', 'md', 'json', 'xml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
                
        elif ext == 'csv':
            df = pd.read_csv(file_path)
            text = df.to_markdown(index=False)
            
        elif ext in ['xlsx', 'xls']:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            for sheet, df in df_dict.items():
                text += f"\n\n### 表格：{sheet}\n\n"
                # fillna("") 防止 NaN 破坏 Markdown 渲染
                text += df.fillna("").astype(str).to_markdown(index=False)
                
        elif ext in ['docx']:
            doc = docx.Document(file_path)
            # 1. 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            text += "\n"
            # 2. 提取表格并转为 Markdown 格式
            for table in doc.tables:
                for i, row in enumerate(table.rows):
                    row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                    text += "| " + " | ".join(row_data) + " |\n"
                    if i == 0:
                        text += "|" + "|".join(["---"] * len(row.cells)) + "|\n"
                text += "\n"
                
        elif ext in ['pptx']:
            prs = pptx.Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text += f"\n\n\n\n"
                # 1. 提取文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    # 2. 提取幻灯片里的表格
                    if shape.has_table:
                        for r_idx, row in enumerate(shape.table.rows):
                            row_data = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
                            text += "| " + " | ".join(row_data) + " |\n"
                            if r_idx == 0:
                                text += "|" + "|".join(["---"] * len(row.cells)) + "|\n"
                        text += "\n"
    except Exception as e:
        raise RuntimeError(f"解析 {ext.upper()} 文件失败: {str(e)}")
        
    return text